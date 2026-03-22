import os
import datetime
import streamlit as st
from openai import OpenAI
import PyPDF2
import docx
import pptx
import streamlit.components.v1 as components

# 设置页面宽屏显示
st.set_page_config(page_title="智能学习助手", page_icon="🎓", layout="centered")

# --- 1. 安全获取 API 密钥 (云端/本地双重适配) ---
try:
    # 优先尝试从 Streamlit 云端的机密设置中读取
    api_key = st.secrets["DEEPSEEK_API_KEY"]
except:
    # 如果你在本地自己电脑上测试，它会回退使用下面这个备用密钥
    # ⚠️ 警告：分享给别人或上传 GitHub 前，请把下面这行清空或填假字母！
    api_key = "你的DeepSeek真实密钥"

os.environ["DEEPSEEK_API_KEY"] = api_key

st.title("🎓 引导式智能学习助手")


@st.cache_resource
def get_client():
    return OpenAI(
        api_key=os.environ.get("DEEPSEEK_API_KEY"),
        base_url="https://api.deepseek.com"
    )


client = get_client()
current_time = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")

system_instruction = f"""
1. 角色与目标
你是一位教育心理学专家。你的核心目标是执行“无限期引导式学习”，严禁直接向用户灌输最终结论。

2. 核心规矩（最高优先级 - 绝对防线）
- 🚫 强制两轮底线：在用户针对你的引导提问进行至少两轮实质性的思考与回复之前，如果用户企图走捷径（如输入“给出答案”），你绝对禁止给出答案！你必须温和但坚定地拒绝，并把话题拉回。
- 听从“触发指令”：只有在满足了上述前提下，且用户明确输入“请给出完整答案”时，你才可以进行完整的推导和解答。
- 绝不擅自结束：即使给出了完整答案，你也必须询问“对于以上推导，你还有哪里不清楚吗？”，等待追问。
- 终极指令：只有当用户明确输入“结束本次学习”时，你才进行本轮知识点的全面大总结，并正式宣布指导结束。

3. 情境感知
当前时间是 {current_time}。在引导时，可以适时结合现实情境来举例。

4. 输出格式与约束
- 绝对禁止使用 \\(、\\) 或 \\[、\\] 来包裹数学公式。必须使用标准的美元符号：行内单美元 `$ $`，独立块双美元 `$$ $$`。
- 如果用户上传了参考资料，请务必仔细阅读，并在提问中紧密结合资料核心概念。
"""

# --- 状态初始化 ---
if "history" not in st.session_state:
    st.session_state.history = []
if "viewing_past" not in st.session_state:
    st.session_state.viewing_past = None
if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "system", "content": system_instruction}]


# --- 文件提取函数 ---
def extract_text_from_file(uploaded_file):
    file_name = uploaded_file.name.lower()
    text = ""
    try:
        if file_name.endswith('.txt'):
            text = uploaded_file.getvalue().decode("utf-8")
        elif file_name.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_name.endswith('.docx'):
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_name.endswith('.pptx'):
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        return f"读取文件失败: {e}"
    return text


# --- 核心逻辑：开启新会话 ---
def start_new_chat():
    if len(st.session_state.messages) > 1:
        title = "新学习主题"
        for m in st.session_state.messages:
            if m["role"] == "user" and not m["content"].startswith("[系统提示："):
                title = m["content"][:12] + "..."
                break
        st.session_state.history.append({
            "title": title,
            "messages": st.session_state.messages.copy()
        })
    st.session_state.messages = [{"role": "system", "content": system_instruction}]
    st.session_state.viewing_past = None


# ================= 侧边栏 UI =================
with st.sidebar:
    if st.button("➕ 开启新一轮学习", use_container_width=True, type="primary"):
        start_new_chat()
        st.rerun()

    st.divider()

    st.header("🛠️ 学习资料库")
    uploaded_file = st.file_uploader("上传学习资料 (TXT/PDF/Word/PPT)", type=["txt", "pdf", "docx", "pptx"])
    if uploaded_file is not None:
        with st.spinner('正在努力读取文件内容...'):
            file_content = extract_text_from_file(uploaded_file)
            if "读取文件失败" not in file_content and file_content.strip() != "":
                st.success(f"成功读取：{uploaded_file.name}！")
                if not any(f"读取了资料：{uploaded_file.name}" in m["content"] for m in st.session_state.messages):
                    st.session_state.messages.append({
                        "role": "user",
                        "content": f"[系统提示：用户刚上传并读取了资料：{uploaded_file.name}。请作为核心参考背景]\n\n{file_content}"
                    })
            else:
                st.error("无法提取文字，可能是纯扫描件或图片。")

    st.divider()

    st.subheader("📚 历史学习记录")
    if not st.session_state.history:
        st.caption("暂无归档记录，结束学习后点击上方按钮即可归档。")
    else:
        for i, past_session in enumerate(st.session_state.history):
            if st.button(f"📝 {past_session['title']}", key=f"hist_{i}", use_container_width=True):
                st.session_state.viewing_past = i
                st.rerun()

    if st.session_state.viewing_past is not None:
        st.info("👀 您当前正在查看历史记录")
        if st.button("🔙 返回当前学习进度", use_container_width=True):
            st.session_state.viewing_past = None
            st.rerun()

    st.divider()

    st.subheader("💾 导出为 PDF")

    current_display_messages = st.session_state.messages if st.session_state.viewing_past is None else \
    st.session_state.history[st.session_state.viewing_past]["messages"]

    if len(current_display_messages) > 1:
        pdf_button_html = """
        <button onclick="window.parent.print()" style="width: 100%; padding: 10px; background-color: #FF4B4B; color: white; border: none; border-radius: 8px; font-size: 16px; font-weight: bold; cursor: pointer; transition: 0.3s;" onmouseover="this.style.backgroundColor='#FF6B6B'" onmouseout="this.style.backgroundColor='#FF4B4B'">
            📄 一键排版并另存为 PDF
        </button>
        """
        components.html(pdf_button_html, height=50)
        st.caption("💡 提示：点击将导出**当前屏幕**上显示的笔记。")

# ================= 主界面 UI =================
display_messages = st.session_state.messages
if st.session_state.viewing_past is not None:
    display_messages = st.session_state.history[st.session_state.viewing_past]["messages"]

for msg in display_messages:
    if msg["role"] != "system" and not msg["content"].startswith("[系统提示："):
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

if st.session_state.viewing_past is None:
    if prompt := st.chat_input("向导师提问，或输入'结束本次学习'..."):
        with st.chat_message("user"):
            st.markdown(prompt)
        st.session_state.messages.append({"role": "user", "content": prompt})

        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            full_response = ""

            responses = client.chat.completions.create(
                model="deepseek-chat",
                messages=st.session_state.messages,
                stream=True,
                temperature=0.7
            )

            for chunk in responses:
                if chunk.choices[0].delta.content is not None:
                    full_response += chunk.choices[0].delta.content
                    message_placeholder.markdown(full_response + "▌")

            message_placeholder.markdown(full_response)

        st.session_state.messages.append({"role": "assistant", "content": full_response})
else:
    st.warning("🔒 历史记录属于只读模式。若要继续学习，请点击左侧栏的【🔙 返回当前学习进度】或【➕ 开启新一轮学习】。")