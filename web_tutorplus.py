import os
import json
import base64
import datetime
import streamlit as st
from openai import OpenAI
import PyPDF2
import docx
import pptx
import streamlit.components.v1 as components

# 设置页面宽屏显示
st.set_page_config(page_title="Web Tutor Plus", page_icon="🎓", layout="wide")

# ==================== 1. 中转站专属配置 ====================
try:
    # 注意：这里我又改回了 OPENAI_API_KEY，记得去后台换成你中转站的钥匙
    api_key = st.secrets["OPENAI_API_KEY"]
except:
    api_key = "sk-..."

if not api_key or api_key.startswith("AIza") or api_key.startswith("sk-..."):
    st.error("⚠️ 请在 Streamlit 后台配置中转站的 OPENAI_API_KEY。")
    st.stop()

# 🚀 重新连回柏拉图中转站
client = OpenAI(
    api_key=api_key,
    base_url="https://api.bltcy.ai/v1"
)

# 🧠 切换模型：你可以用中转站里的 gemini-1.5-pro，或者直接上国内看图最强王者 qwen-vl-max
VISION_MODEL = "qwen-vl-max"

st.title("🎓 Web Tutor Plus")


@st.cache_resource
def get_current_time():
    return datetime.datetime.now().strftime("%Y-%m-%d %H:%M")


# ==================== 新增：自动存档与读档引擎 ====================
BACKUP_FILE = "chat_backup.json"


def load_backup():
    if os.path.exists(BACKUP_FILE):
        try:
            with open(BACKUP_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return None
    return None


def save_backup(messages):
    with open(BACKUP_FILE, "w", encoding="utf-8") as f:
        json.dump(messages, f, ensure_ascii=False, indent=2)


# ==================================================================

# ==================== 2. 全面升级的系统提示词 ====================
system_instruction = f"""
1. 角色与目标
你是一位顶级的教育心理学专家。你的核心目标是执行“无限期引导式学习”，严禁直接向用户灌输最终结论。你现在具备顶级的视觉能力，可以看懂最复杂的数学公式、物理受力图和工程代码。

2. 核心规矩（最高优先级 - 绝对防线）
- 🚫 强制两轮底线：在用户针对你的引导提问进行至少两轮实质性的思考与回复之前，如果用户企图走捷径（如输入“给出答案”、“不想猜”、“直接告诉我”等），你绝对禁止给出答案！你必须温和但坚定地拒绝，并把话题拉回。
- 听从“触发指令”：只有在满足了上述前提下，且用户明确输入“请给出完整答案”时，你才可以进行完整的推导和解答。
- 绝不擅自结束：即使给出了完整答案，你也必须询问“对于以上推导，你还有哪里不清楚吗？”，等待追问。
- 终极指令：只有当用户明确输入“结束本次学习”时，你才进行总结，并正式宣布指导结束。

3. 图片理解规则 (Vision)
- 如果用户上传了图片（例如常微分方程推导截图、经典力学非惯性系受力图、CNC加工图纸等），你必须优先、极其仔细地阅读并分析图片内容。
- 在后续的引导和提问中，要紧密结合图片中的细节（如特定的变量符号、受力方向、微积分上下限）。

4. 文档资料库规则
- 如果用户上传了文档，仔细阅读并在提问中结合资料核心概念。

5. 输出格式与约束
- 数学公式：美元符号：行内 `$ $`，独立块双美元 `$$ $$`。
"""

# --- 状态初始化 ---
if "history" not in st.session_state:
    st.session_state.history = []
if "viewing_past" not in st.session_state:
    st.session_state.viewing_past = None

if "messages" not in st.session_state:
    backup = load_backup()
    if backup:
        st.session_state.messages = backup
    else:
        st.session_state.messages = [{"role": "system", "content": [{"type": "text", "text": system_instruction}]}]


# --- 实用辅助函数 ---
def encode_image(image_file):
    return base64.b64encode(image_file.getvalue()).decode('utf-8')


def extract_text_from_file(uploaded_file):
    file_name = uploaded_file.name.lower()
    text = ""
    try:
        if file_name.endswith('.txt'):
            text = uploaded_file.getvalue().decode("utf-8")
        elif file_name.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        elif file_name.endswith('.docx'):
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_name.endswith('.pptx'):
            ppt = pptx.Presentation(uploaded_file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
    except Exception as e:
        return f"读取文件失败: {e}"
    return text


def start_new_chat():
    if len(st.session_state.messages) > 1:
        title = "新学习主题"
        for m in st.session_state.messages:
            if m["role"] == "user":
                if isinstance(m["content"], list):
                    for item in m["content"]:
                        if item["type"] == "text":
                            title = item["text"][:12] + "..."
                            break
                    if title != "新学习主题": break
                else:
                    title = m["content"][:12] + "..."
                    break

        st.session_state.history.append({
            "title": title,
            "messages": st.session_state.messages.copy()
        })

    st.session_state.messages = [{"role": "system", "content": [{"type": "text", "text": system_instruction}]}]
    st.session_state.viewing_past = None
    save_backup(st.session_state.messages)


# ==================== 3. 侧边栏 UI ====================
with st.sidebar:
    st.button("➕ 开启新一轮学习", use_container_width=True, type="primary", on_click=start_new_chat)
    st.divider()

    st.header("🛠️ 学习资料库")

    st.subheader("🖼️ 上传图片理解")
    uploaded_image = st.file_uploader("上传公式、受力图、代码截图", type=["png", "jpg", "jpeg"])
    if uploaded_image is not None:
        if st.session_state.viewing_past is None:
            with st.spinner('正在使用视觉模型分析图片...'):
                base64_image = encode_image(uploaded_image)
                already_uploaded = False
                for msg in st.session_state.messages:
                    if msg["role"] == "user" and isinstance(msg["content"], list):
                        for item in msg["content"]:
                            if item["type"] == "image_url" and uploaded_image.name in msg.get("metadata", {}).get(
                                    "file_name", ""):
                                already_uploaded = True
                                break
                    if already_uploaded: break

                if not already_uploaded:
                    st.session_state.messages.append({
                        "role": "user",
                        "content": [
                            {"type": "text",
                             "text": f"[用户上传并分析了图片：{uploaded_image.name}，请仔细查看图片内容并将其融入你的引导过程]"},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                        ],
                        "metadata": {"file_name": uploaded_image.name}
                    })
                    st.success(f"成功分析：{uploaded_image.name}！")

    st.divider()

    st.subheader("📄 上传文档资料")
    uploaded_file = st.file_uploader("上传TXT/PDF/Word/PPT", type=["txt", "pdf", "docx", "pptx"])
    if uploaded_file is not None:
        if st.session_state.viewing_past is None:
            with st.spinner('正在读取文件内容...'):
                file_content = extract_text_from_file(uploaded_file)
                if "读取文件失败" not in file_content and file_content.strip() != "":
                    if not any(f"读取了资料：{uploaded_file.name}" in m["content"][0]["text"] if isinstance(
                            m.get("content", []), list) else f"读取了资料：{uploaded_file.name}" in m.get("content", "")
                               for m in st.session_state.messages):
                        st.session_state.messages.append({
                            "role": "user",
                            "content": [{"type": "text",
                                         "text": f"[系统提示：用户刚上传并读取了资料：{uploaded_file.name}。以下是资料内容，请在后续对话中作为参考背景]\n\n{file_content}"}]
                        })
                        st.success(f"成功读取：{uploaded_file.name}！")
                else:
                    st.error("无法提取文字，可能是纯扫描件或图片。")

    st.divider()

    st.subheader("📚 历史学习记录")
    if not st.session_state.history:
        st.caption("暂无归档记录，结束学习后点击上方按钮即可归档。")
    else:
        for i, past_session in enumerate(st.session_state.history):
            if st.button(f"📝 {past_session['title']}", key=f"hist_{i}", use_container_width=True):
                st.session_state.viewing_past = i
                st.rerun()

    if st.session_state.viewing_past is not None:
        st.info("👀 您当前正在查看历史记录")
        if st.button("🔙 返回当前学习进度", use_container_width=True):
            st.session_state.viewing_past = None
            st.rerun()

    st.divider()

    st.subheader("💾 导出为 PDF")
    current_display_messages = st.session_state.messages if st.session_state.viewing_past is None else \
    st.session_state.history[st.session_state.viewing_past]["messages"]

    if len(current_display_messages) > 1:
        # 💥 这里已经替换为最新的【多重防爆版】PDF引擎
        pdf_button_html = """
            <script src="https://cdnjs.cloudflare.com/ajax/libs/html2pdf.js/0.10.1/html2pdf.bundle.min.js"></script>

            <button id="pdf-btn" onclick="generatePDF()" style="width: 100%; padding: 10px; background-color: #FF4B4B; color: white; border: none; border-radius: 8px; font-size: 16px; font-weight: bold; cursor: pointer; transition: 0.3s; box-shadow: 0 2px 5px rgba(0,0,0,0.2);">
                📄 下载纯净笔记 PDF
            </button>

            <script>
        function generatePDF() {
            var btn = document.getElementById('pdf-btn');
            btn.innerText = '⏳ 正在排版纯净笔记...';
            btn.style.backgroundColor = '#ff8b8b';

            try {
                // 1. 【防弹抓取】获取包含所有对话的最里层容器，避开外层滚动条的干扰
                var mainElement = window.parent.document.querySelector('[data-testid="stMainBlockContainer"]') || 
                                  window.parent.document.querySelector('.main') || 
                                  window.parent.document.body;
                
                if (!mainElement) throw new Error("找不到主界面元素");
                var clonedMain = mainElement.cloneNode(true);

                // 💥【终极黑魔法】专治白屏：强制展开所有隐藏/滚动区域，铺平内容
                clonedMain.style.height = 'auto';
                clonedMain.style.overflow = 'visible';
                clonedMain.style.position = 'relative';

                // 2. 安全清除垃圾元素
                var safeRemove = function(selector) {
                    var els = clonedMain.querySelectorAll(selector);
                    els.forEach(function(e) { e.remove(); });
                };

                safeRemove('h1'); 
                safeRemove('div[data-testid="stAlert"]'); 
                safeRemove('[data-testid="stChatAvatar"]'); 

                try {
                    var allMessages = clonedMain.querySelectorAll('div[data-testid="stChatMessage"]');
                    allMessages.forEach(function(msg) {
                        if (msg.querySelector('div[aria-label="user"]')) {
                            msg.remove();
                        }
                    });
                } catch (e) { console.log("清理跳过"); }

                // 3. 💥 PDF引擎防白屏配置
                var opt = {
                  margin:       [0.4, 0.4, 0.4, 0.4],
                  filename:     'WebTutor_学习笔记.pdf',
                  image:        { type: 'jpeg', quality: 0.98 },
                  html2canvas:  { 
                      scale: 2, 
                      useCORS: true, 
                      logging: false,
                      scrollY: 0,  // 强制视角回到最顶部，防止内容向下偏移
                      windowHeight: mainElement.scrollHeight // 强制引擎读取网页的真实物理总高度
                  },
                  jsPDF:        { unit: 'in', format: 'a4', orientation: 'portrait' }
                };

                // 4. 强制输出
                html2pdf().set(opt).from(clonedMain).save().then(function() {
                    btn.innerText = '📄 下载纯净笔记 PDF';
                    btn.style.backgroundColor = '#FF4B4B';
                }).catch(function(err) {
                    btn.innerText = '❌ 导出错误，点击重试';
                    btn.style.backgroundColor = '#FF4B4B';
                });
                
            } catch (error) {
                btn.innerText = '⚠️ 页面提取失败';
                btn.style.backgroundColor = '#FF4B4B';
                setTimeout(() => { btn.innerText = '📄 下载纯净笔记 PDF'; }, 3000);
            }
        }
        </script>
            """
        components.html(pdf_button_html, height=80)

# ==================== 4. 主界面 UI ====================
display_messages = st.session_state.messages
if st.session_state.viewing_past is not None:
    display_messages = st.session_state.history[st.session_state.viewing_past]["messages"]

for msg in display_messages:
    if msg["role"] == "system": continue

    with st.chat_message(msg["role"]):
        if isinstance(msg["content"], list):
            text_content = ""
            for item in msg["content"]:
                if item["type"] == "text":
                    if item["text"].startswith("[系统提示：") or item["text"].startswith(
                        "[用户上传并分析了图片："): continue
                    text_content += item["text"]
                elif item["type"] == "image_url":
                    st.image(item["image_url"]["url"], caption="您上传的图片资料", width=400)

            if text_content:
                st.markdown(text_content)
        else:
            if not msg["content"].startswith("[系统提示："):
                st.markdown(msg["content"])

# 聊天输入与逻辑处理
if st.session_state.viewing_past is None:
    if prompt := st.chat_input("向导师提问，或输入'结束本次学习'以归档..."):
        with st.chat_message("user"):
            st.markdown(prompt)
        st.session_state.messages.append({"role": "user", "content": prompt})
        save_backup(st.session_state.messages)

        with st.chat_message("assistant"):
            message_placeholder = st.empty()
            full_response = ""

            try:
                responses = client.chat.completions.create(
                    model=VISION_MODEL,
                    messages=st.session_state.messages,
                    stream=True,
                    temperature=0.7
                )

                for chunk in responses:
                    if hasattr(chunk, 'choices') and chunk.choices and len(chunk.choices) > 0:
                        content = getattr(chunk.choices[0].delta, 'content', None)
                        if content is not None:
                            full_response += content
                            message_placeholder.markdown(full_response + "▌")

            except Exception as e:
                st.error(f"请求出错 (请检查密钥或网络): {e}")

            message_placeholder.markdown(full_response)
            st.session_state.messages.append({"role": "assistant", "content": full_response})
            save_backup(st.session_state.messages)

else:
    st.warning("🔒 历史记录属于只读模式。若要继续学习，请点击左侧栏的【🔙 返回当前学习进度】或【➕ 开启新一轮学习】。")
